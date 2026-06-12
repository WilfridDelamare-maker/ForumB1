from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Couleurs
DARK_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # fond sombre
ACCENT      = RGBColor(0x89, 0xB4, 0xFA)   # bleu clair
ACCENT2     = RGBColor(0xA6, 0xE3, 0xA1)   # vert clair
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD0, 0xDA)
DARK_CARD   = RGBColor(0x31, 0x32, 0x44)   # fond carte
PINK        = RGBColor(0xF3, 0x8B, 0xA8)
YELLOW      = RGBColor(0xF9, 0xE2, 0xAF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# ─── helpers ─────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide

def txt(slide, text, left, top, width, height,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return tb

def rect(slide, left, top, width, height, fill_color, radius=False):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def card(slide, left, top, width, height):
    return rect(slide, left, top, width, height, DARK_CARD)

def accent_bar(slide, left, top, width=0.06, height=0.55):
    return rect(slide, left, top, width, height, ACCENT)

def bullet_block(slide, items, left, top, width, size=17, color=LIGHT_GRAY, spacing=0.38):
    for i, item in enumerate(items):
        txt(slide, "▸  " + item, left, top + i * spacing, width, spacing + 0.1,
            size=size, color=color)

def section_title(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.35, DARK_CARD)
    accent_bar(slide, 0.35, 0.38)
    txt(slide, title, 0.55, 0.32, 11, 0.75, size=30, bold=True, color=ACCENT)
    if subtitle:
        txt(slide, subtitle, 0.55, 0.95, 11, 0.45, size=15, color=LIGHT_GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()

rect(slide, 0, 2.6, 13.33, 0.08, ACCENT)

txt(slide, "FORUM B1", 0, 0.9, 13.33, 1.4,
    size=70, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, "Application web de forum — Go · SQLite · Docker",
    0, 2.1, 13.33, 0.7, size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

infos = [
    ("Langage",    "Go 1.25"),
    ("BDD",        "SQLite"),
    ("Auth",       "Local + GitHub + Google OAuth"),
    ("Deploy",     "Docker / docker-compose"),
]
card_w, gap = 2.8, 0.18
start_x = (13.33 - (len(infos) * card_w + (len(infos)-1) * gap)) / 2
for i, (label, val) in enumerate(infos):
    cx = start_x + i * (card_w + gap)
    card(slide, cx, 3.1, card_w, 1.4)
    txt(slide, label, cx, 3.2, card_w, 0.4, size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, val,   cx, 3.65, card_w, 0.65, size=16, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "Projet réalisé en équipe — B1", 0, 6.8, 13.33, 0.5,
    size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Présentation du projet
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Présentation du projet", "Un forum complet développé from scratch en Go")

txt(slide, "Qu'est-ce que c'est ?", 0.5, 1.55, 6, 0.5, size=19, bold=True, color=ACCENT)
bullet_block(slide, [
    "Forum web multithématique (posts, commentaires, réactions)",
    "Inscription, connexion, déconnexion sécurisées",
    "Connexion via compte GitHub ou Google (OAuth 2.0)",
    "Catégorisation des posts, recherche par mot-clé",
    "Système de likes / dislikes sur posts et commentaires",
    "Upload d'images dans les posts (JPEG, PNG, GIF)",
    "Modification et suppression de ses propres posts",
], 0.5, 2.05, 6.1, size=16)

txt(slide, "Pourquoi Go ?", 7.2, 1.55, 5.5, 0.5, size=19, bold=True, color=ACCENT2)
bullet_block(slide, [
    "Serveur HTTP natif (net/http) — pas de framework",
    "Performances élevées, faible consommation mémoire",
    "Compilation statique → déploiement simple",
    "Typage fort, gestion explicite des erreurs",
    "Idéal pour apprendre les bases du web backend",
], 7.2, 2.05, 5.8, size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Architecture du projet", "Organisation des packages Go")

packages = [
    ("handlers/",  ACCENT,  "Reçoit les requêtes HTTP\nGère la logique métier\nAppelle database/"),
    ("database/",  ACCENT2, "Requêtes SQL SQLite\nCRUD complet\nTransactions"),
    ("models/",    YELLOW,  "Structures de données\nPost, User, Comment…\nTemplateData"),
    ("templates/", PINK,    "Fichiers .tmpl (HTML)\nRendu côté serveur\nDonnées injectées"),
    ("fake/",      LIGHT_GRAY, "Couche d'abstraction\nentre handlers\net database"),
]

col_w = 2.3
start = 0.45
for i, (name, color, desc) in enumerate(packages):
    cx = start + i * (col_w + 0.12)
    card(slide, cx, 1.6, col_w, 4.5)
    rect(slide, cx, 1.6, col_w, 0.06, color)
    txt(slide, name, cx, 1.7, col_w, 0.55, size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, desc, cx, 2.35, col_w, 3.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(slide, "main.go  →  déclare toutes les routes HTTP et démarre le serveur sur :8085",
    0.5, 6.35, 12.3, 0.5, size=15, color=WHITE, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Authentification
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Authentification", "3 méthodes de connexion — sessions cookie HttpOnly")

# Colonne locale
card(slide, 0.4, 1.55, 3.7, 4.9)
txt(slide, "Locale", 0.4, 1.65, 3.7, 0.5, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
bullet_block(slide, [
    "Email + mot de passe",
    "Hachage bcrypt",
    "Email unique en BDD",
    "Session UUID 24h",
    "Cookie HttpOnly",
], 0.65, 2.25, 3.2, size=15, spacing=0.42)

# GitHub
card(slide, 4.6, 1.55, 3.7, 4.9)
txt(slide, "GitHub OAuth 2.0", 4.6, 1.65, 3.7, 0.5, size=18, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
bullet_block(slide, [
    "Redirection → github.com",
    "Callback /auth/github",
    "Récupère email + login",
    "Crée compte si nouveau",
    "provider = 'github'",
], 4.85, 2.25, 3.2, size=15, color=LIGHT_GRAY, spacing=0.42)

# Google
card(slide, 8.8, 1.55, 3.7, 4.9)
txt(slide, "Google OAuth 2.0", 8.8, 1.65, 3.7, 0.5, size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
bullet_block(slide, [
    "Redirection → google.com",
    "Callback /auth/google",
    "Récupère email + nom",
    "Crée compte si nouveau",
    "provider = 'google'",
], 9.05, 2.25, 3.2, size=15, color=LIGHT_GRAY, spacing=0.42)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Fonctionnalités posts & commentaires
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Posts & Commentaires", "Cœur du forum")

# Posts
card(slide, 0.4, 1.55, 5.9, 5.1)
txt(slide, "Posts", 0.6, 1.7, 5.5, 0.55, size=22, bold=True, color=ACCENT)
bullet_block(slide, [
    "Titre, contenu, 1 à N catégories",
    "Image optionnelle (max 20 Mo)",
    "Validation du type par magic bytes",
    "Auteur stocké → seul lui peut modifier",
    "Modification titre + contenu",
    "Suppression définitive",
    "Tri du plus récent au plus ancien",
], 0.6, 2.3, 5.5, size=16, spacing=0.4)

# Commentaires
card(slide, 6.9, 1.55, 5.9, 5.1)
txt(slide, "Commentaires", 7.1, 1.7, 5.5, 0.55, size=22, bold=True, color=ACCENT2)
bullet_block(slide, [
    "Liés à un post existant",
    "Contenu non vide requis",
    "Connexion obligatoire",
    "Likes / Dislikes sur commentaires",
    "Affichage du nombre de réactions",
    "Auteur affiché avec chaque commentaire",
    "Horodatage automatique",
], 7.1, 2.3, 5.5, size=16, color=LIGHT_GRAY, spacing=0.4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Base de données
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Base de données", "SQLite — 6 tables — Transactions")

tables = [
    ("users",           ["id", "email", "username", "password_hash", "provider", "provider_id", "created_at"]),
    ("sessions",        ["id (UUID)", "user_id → users", "expires_at"]),
    ("categories",      ["id", "name (unique)"]),
    ("posts",           ["id", "title", "content", "image_path", "author_id → users", "created_at"]),
    ("post_categories", ["post_id → posts", "category_id → categories", "(PK composite)"]),
    ("likes",           ["id", "user_id → users", "post_id?", "comment_id?", "value (+1/-1)"]),
]

col_w, col_h = 3.8, 2.2
positions = [
    (0.3, 1.5), (4.4, 1.5), (8.5, 1.5),
    (0.3, 4.0), (4.4, 4.0), (8.5, 4.0),
]
colors_t = [ACCENT, ACCENT2, YELLOW, PINK, LIGHT_GRAY, ACCENT]

for (lx, ly), (name, fields), color in zip(positions, tables, colors_t):
    card(slide, lx, ly, col_w, col_h)
    rect(slide, lx, ly, col_w, 0.045, color)
    txt(slide, name, lx, ly + 0.05, col_w, 0.4, size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    field_txt = "\n".join("  · " + f for f in fields)
    txt(slide, field_txt, lx + 0.1, ly + 0.5, col_w - 0.2, col_h - 0.6, size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Routes HTTP
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Routes HTTP", "net/http standard library — mux personnalisé")

routes = [
    ("GET",  "/",                      "Accueil — liste posts + recherche",           ACCENT),
    ("GET",  "/register  /login",      "Formulaires inscription / connexion",          ACCENT),
    ("POST", "/register  /login",      "Traitement formulaires + création session",    ACCENT2),
    ("POST", "/logout",                "Suppression cookie + session BDD",             ACCENT2),
    ("GET",  "/posts/{id}",            "Affiche un post et ses commentaires",          YELLOW),
    ("POST", "/posts/{id}/comments",   "Ajoute un commentaire (auth requis)",          YELLOW),
    ("GET",  "/posts/create",          "Formulaire création post (auth requis)",       PINK),
    ("POST", "/posts/create",          "Enregistre post + image (auth requis)",        PINK),
    ("POST", "/posts/{id}/like|dislike","Like ou dislike sur un post",                 LIGHT_GRAY),
    ("GET",  "/posts/{id}/edit",       "Formulaire édition (auteur uniquement)",       LIGHT_GRAY),
    ("POST", "/posts/{id}/edit|delete","Modifie ou supprime le post",                  LIGHT_GRAY),
    ("GET",  "/categories  /random",   "Liste catégories / post aléatoire",            ACCENT),
    ("GET",  "/auth/github|google",    "Redirection OAuth",                            ACCENT2),
]

row_h = 0.37
col_method = 0.35
col_path   = 1.3
col_desc   = 5.3
for i, (method, path, desc, color) in enumerate(routes):
    y = 1.55 + i * row_h
    if i % 2 == 0:
        rect(slide, 0.3, y, 12.73, row_h, RGBColor(0x28, 0x29, 0x3D))
    txt(slide, method, col_method, y, 0.85, row_h, size=13, bold=True, color=color)
    txt(slide, path,   col_path,   y, 3.7,  row_h, size=13, color=WHITE)
    txt(slide, desc,   col_desc,   y, 7.8,  row_h, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Déploiement Docker
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Déploiement", "Docker + docker-compose — build multi-stage")

# Container db
card(slide, 0.4, 1.6, 5.8, 2.3)
txt(slide, "Container  db", 0.6, 1.7, 5.4, 0.5, size=17, bold=True, color=ACCENT)
bullet_block(slide, [
    "Image : alpine:3.19 (légère)",
    "Crée /data et reste en vie",
    "Healthcheck avant démarrage forum",
    "Volume forum_data partagé",
], 0.6, 2.2, 5.4, size=15, spacing=0.37)

# Container forum
card(slide, 6.7, 1.6, 6.2, 2.3)
txt(slide, "Container  forum", 6.9, 1.7, 5.8, 0.5, size=17, bold=True, color=ACCENT2)
bullet_block(slide, [
    "Build multi-stage (builder + runtime)",
    "CGO activé pour go-sqlite3",
    "Binaire statique — image finale Debian slim",
    "Port 8085 exposé",
], 6.9, 2.2, 5.8, size=15, color=LIGHT_GRAY, spacing=0.37)

# Multi-stage detail
card(slide, 0.4, 4.15, 12.5, 2.0)
txt(slide, "Build multi-stage", 0.6, 4.25, 12, 0.45, size=17, bold=True, color=YELLOW)
bullet_block(slide, [
    "Stage 1 (builder) : golang:1.26-alpine + gcc → compile le binaire avec CGO_ENABLED=1",
    "Stage 2 (runtime) : debian:bookworm-slim → copie uniquement le binaire + templates + static",
    "Résultat : image finale légère (~30 Mo), sans toolchain Go, démarrage rapide",
], 0.6, 4.7, 12, size=15, color=LIGHT_GRAY, spacing=0.42)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Points techniques notables
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Points techniques notables", "Choix d'implémentation intéressants")

points = [
    (ACCENT,     "Magic bytes pour validation d'image",
                 "Le type MIME n'est pas vérifié via l'extension (falsifiable)\nmais par lecture des premiers octets du fichier (JPEG: FF D8 FF, PNG: 89 50 4E 47, GIF: 47 49 46)"),
    (ACCENT2,    "Requête SQL partagée (postSelectQuery)",
                 "Une constante SQL de base est réutilisée par GetAllPosts, GetPostByID et GetPostsByCategory\npour éviter la duplication. Chaque fonction ajoute juste sa clause WHERE/GROUP BY."),
    (YELLOW,     "Interface scanPost pour Rows et Row",
                 "Une seule fonction scanPost accepte une interface { Scan(...any) error }\nce qui évite de dupliquer le scan pour *sql.Rows (Query) et *sql.Row (QueryRow)."),
    (PINK,       "Sessions côté serveur",
                 "L'UUID de session est stocké en BDD avec une expiration 24h.\nLe cookie HttpOnly empêche l'accès JavaScript → protection XSS."),
]

for i, (color, title, desc) in enumerate(points):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.55
    ly = 1.6 + row * 2.5
    card(slide, lx, ly, 6.25, 2.25)
    rect(slide, lx, ly, 0.05, 2.25, color)
    txt(slide, title, lx + 0.2, ly + 0.1, 5.9, 0.5, size=16, bold=True, color=color)
    txt(slide, desc,  lx + 0.2, ly + 0.65, 5.9, 1.5, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Bilan & démo
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
section_title(slide, "Bilan & Démo", "Ce qu'on a appris, ce qui reste à faire")

card(slide, 0.4, 1.6, 5.8, 5.1)
txt(slide, "Ce que le projet couvre", 0.6, 1.7, 5.5, 0.5, size=18, bold=True, color=ACCENT2)
bullet_block(slide, [
    "Serveur HTTP from scratch en Go",
    "CRUD complet sur 3 entités",
    "OAuth 2.0 réel (GitHub + Google)",
    "Sécurité : bcrypt, HttpOnly, transactions",
    "Upload de fichiers avec validation",
    "Containerisation Docker multi-stage",
    "Routing avancé (Go 1.22 patterns)",
], 0.6, 2.25, 5.5, size=16, color=LIGHT_GRAY, spacing=0.4)

card(slide, 6.9, 1.6, 5.8, 5.1)
txt(slide, "Améliorations possibles", 7.1, 1.7, 5.5, 0.5, size=18, bold=True, color=PINK)
bullet_block(slide, [
    "Pagination des posts",
    "Profil utilisateur modifiable",
    "Notifications en temps réel (WebSocket)",
    "Tests unitaires et d'intégration",
    "Middleware d'authentification centralisé",
    "Rate limiting sur les formulaires",
    "CI/CD GitHub Actions",
], 7.1, 2.25, 5.5, size=16, color=LIGHT_GRAY, spacing=0.4)

# ─── Save ────────────────────────────────────────────────────────────────────
output = "ForumB1_Presentation.pptx"
prs.save(output)
print(f"Fichier généré : {output}")
